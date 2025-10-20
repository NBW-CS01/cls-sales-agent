#!/usr/bin/env python3
"""
Convert Jamie proposal text to PowerPoint presentation
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import re
import sys


def parse_proposal(text):
    """Parse markdown proposal into structured slides"""
    slides = []
    current_slide = None

    lines = text.split('\n')

    for line in lines:
        line = line.strip()

        if not line:
            continue

        # Main title (# Title)
        if line.startswith('# '):
            if current_slide:
                slides.append(current_slide)
            current_slide = {
                'type': 'title',
                'title': line[2:],
                'content': []
            }

        # Subtitle (## Subtitle)
        elif line.startswith('## '):
            if current_slide:
                slides.append(current_slide)
            current_slide = {
                'type': 'section',
                'title': line[3:],
                'content': []
            }

        # Section heading (### Heading)
        elif line.startswith('### '):
            if current_slide:
                slides.append(current_slide)
            current_slide = {
                'type': 'content',
                'title': line[4:],
                'content': []
            }

        # Subheading (#### Subheading)
        elif line.startswith('#### '):
            if current_slide:
                current_slide['content'].append({
                    'type': 'subheading',
                    'text': line[5:]
                })

        # Bullet points
        elif line.startswith('- ') or line.startswith('* '):
            if current_slide:
                current_slide['content'].append({
                    'type': 'bullet',
                    'text': line[2:]
                })

        # Table rows
        elif line.startswith('|'):
            if current_slide:
                if not any(item.get('type') == 'table' for item in current_slide['content']):
                    current_slide['content'].append({
                        'type': 'table',
                        'rows': []
                    })
                # Find the table item and add row
                for item in current_slide['content']:
                    if item.get('type') == 'table':
                        row = [cell.strip() for cell in line.split('|')[1:-1]]
                        if not all(cell.startswith('-') for cell in row):  # Skip separator rows
                            item['rows'].append(row)
                        break

        # Regular paragraphs
        else:
            if current_slide and line:
                current_slide['content'].append({
                    'type': 'text',
                    'text': line
                })

    if current_slide:
        slides.append(current_slide)

    return slides


def create_presentation(slides, output_file):
    """Create PowerPoint presentation from parsed slides"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    for slide_data in slides:
        if slide_data['type'] == 'title':
            # Title slide
            slide_layout = prs.slide_layouts[0]  # Title Slide layout
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            title.text = slide_data['title']

            # Add subtitle if available
            if slide_data['content']:
                subtitle = slide.placeholders[1]
                subtitle.text = '\n'.join([
                    item['text'] for item in slide_data['content']
                    if item['type'] == 'text'
                ])

        elif slide_data['type'] == 'section':
            # Section header slide
            slide_layout = prs.slide_layouts[2]  # Section Header layout
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            title.text = slide_data['title']

        else:
            # Content slide
            slide_layout = prs.slide_layouts[1]  # Title and Content layout
            slide = prs.slides.add_slide(slide_layout)

            title = slide.shapes.title
            title.text = slide_data['title']

            # Add content
            body_shape = slide.shapes.placeholders[1]
            tf = body_shape.text_frame
            tf.clear()

            for item in slide_data['content']:
                if item['type'] == 'bullet':
                    p = tf.add_paragraph()
                    p.text = item['text']
                    p.level = 0
                    # Clean up markdown bold
                    p.text = re.sub(r'\*\*(.*?)\*\*', r'\1', p.text)

                elif item['type'] == 'subheading':
                    p = tf.add_paragraph()
                    p.text = item['text']
                    p.level = 0
                    p.font.bold = True
                    p.font.size = Pt(18)

                elif item['type'] == 'text':
                    p = tf.add_paragraph()
                    p.text = item['text']
                    p.level = 0

                elif item['type'] == 'table':
                    # Add table to slide (simplified - just add as text for now)
                    for row in item['rows']:
                        p = tf.add_paragraph()
                        p.text = ' | '.join(row)
                        p.level = 0
                        p.font.size = Pt(12)

    prs.save(output_file)
    print(f"✅ PowerPoint saved to: {output_file}")


def main():
    if len(sys.argv) < 2:
        print("Usage: python3 proposal-to-pptx.py <proposal.txt> [output.pptx]")
        sys.exit(1)

    input_file = sys.argv[1]
    output_file = sys.argv[2] if len(sys.argv) > 2 else input_file.replace('.txt', '.pptx')

    print(f"📄 Reading proposal from: {input_file}")

    with open(input_file, 'r') as f:
        text = f.read()

    print("🔄 Parsing proposal structure...")
    slides = parse_proposal(text)
    print(f"📊 Found {len(slides)} slides")

    print("🎨 Creating PowerPoint presentation...")
    create_presentation(slides, output_file)


if __name__ == '__main__':
    main()
